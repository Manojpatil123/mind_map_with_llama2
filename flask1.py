from flask import Flask, request,send_file
from llama_index import LangchainEmbedding, ServiceContext,StorageContext,load_index_from_storage
from llama_index.llms import LlamaCPP
from llama_index.llms.llama_utils import messages_to_prompt, completion_to_prompt
from llama_index import  SimpleDirectoryReader
from llama_index import GPTVectorStoreIndex
import os
from langchain.embeddings.huggingface import HuggingFaceEmbeddings
import base64
from IPython.display import Image, display
import matplotlib.pyplot as plt
import requests
import re
from langchain.document_loaders import PyPDFLoader
import pypdf
from pptx import Presentation
from docx import Document
from langchain.document_loaders import PyPDFLoader
import pandas as pd



BASE_DIR = os.getcwd()

#loding llm
llm = LlamaCPP(
    # You can pass in the URL to a GGML model to download it automatically
    #model_url="https://huggingface.co/TheBloke/Llama-2-13B-chat-GGML/resolve/main/llama-2-13b-chat.ggmlv3.q4_0.bin",
    # optionally, you can set the path to a pre-downloaded model instead of model_url
    model_path=BASE_DIR+"\models\llama-2-7b-chat.ggmlv3.q4_0.bin",
    temperature=0.1,
    max_new_tokens=1000,
    # llama2 has a context window of 4096 tokens, but we set it lower to allow for some wiggle room
    context_window=4000,
    # kwargs to pass to __call__()
    generate_kwargs={},
    # kwargs to pass to __init__()
    # set to at least 1 to use GPU
    model_kwargs={"n_gpu_layers": 1},
    # transform inputs into Llama2 format
    messages_to_prompt=messages_to_prompt,
    completion_to_prompt=completion_to_prompt,
    verbose=True,
)

#loding embed model
embed_model = LangchainEmbedding(
  HuggingFaceEmbeddings(model_name="sentence-transformers/all-mpnet-base-v2")
)

def pdf(file_path):
    loader = PyPDFLoader(file_path)
    pages = loader.load_and_split()
    document=''
    for i in range(len(pages)):
      document+=pages[i].page_content.replace('\xa0','').replace('\t\r','')
    return document

def textfile(file_path):
    try:
        with open(file_path, 'r') as file:
            file_content = file.read()
            print(file_content)
    except FileNotFoundError:
        print(f"File not found: {file_path}")
    except Exception as e:
        print(f"An error occurred: {e}")
    return file_content

def csvfile(file_path):
    try:
        df = pd.read_csv(file_path)
        for column in df.columns:
            print(df[column])
    except FileNotFoundError:
        print(f"File not found: {file_path}")
    except Exception as e:
        print(f"An error occurred: {e}")
        return df.to_string()

def excelfile(file_path):
    try:
        df = pd.read_excel(file_path)
        for column in df.columns:
            print(df[column])
    except FileNotFoundError:
        print(f"File not found: {file_path}")
    except Exception as e:
        print(f"An error occurred: {e}")
        return df.to_string()

def pptfile(file_path):
    data=''
    try:
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text)
                    data+=shape.text
    except FileNotFoundError:
        print(f"File not found: {file_path}")
    except Exception as e:
        print(f"An error occurred: {e}")
    return data

def docx(file_path):
    try:
        # Load the DOCX file
        doc = Document(file_path)

        # Initialize an empty string to store the extracted text
        extracted_text = ""

        # Iterate through paragraphs in the document and append their text
        for paragraph in doc.paragraphs:
            extracted_text += paragraph.text + "\n"

        return extracted_text
    except Exception as e:
        print(f"An error occurred: {e}")
        return None      

app = Flask(__name__)
'''
central idea: name of topic

* sub topic name
   + summary of sub topic
   * internl sub topic name
       +summary of internal sub topic
* sub topic name
   + summary of sub topic
   + summary of sub topic

give the output of mind map like aabove struture each line proper indentation has to provide, mention each line either * based on sub topic or internel sub topic or + for summary,name of topic should in central idea : name except the output in this format dont generate other texts
'''

@app.route("/mindmap", methods=["POST"])
def upload_files123():
        file = request.files["file"]
        folder_path = BASE_DIR+"/files"
        if not os.path.exists(folder_path):
                os.makedirs(folder_path)
        file_path = os.path.join(folder_path, file.filename)
        print(file.filename)
        file.save(file_path)
        print(file_path)
        if '.pdf' in file.filename:
            documents=pdf(file_path)
        if '.txt' in file.filename:
            documents=textfile(file_path)
        if '.docx' in file.filename:
            documents=docx(file_path)
        if '.pptx' in file.filename:
            documents=pptfile(file_path)
        if '.csv' in file.filename:
            documents=csvfile(file_path)
        if '.xlsx' in file.filename:
            documents=excelfile(file_path)
        
        print(len(documents))
        if len(documents)>10000:
            documents=documents[0:8000]
        '''
        template1="""
you are assistant helps to summarize the given text : {0} and generate output with proper explanation of each section  with paragraph.
        """.format(documents)
        agent1=''
        response_iter = llm.stream_complete(template1)
        for response in response_iter:
           print(response.delta, end="", flush=True)
           agent1+=response.delta'''
        
        template="""
you are assistant helps to create mind map,structure of  Mind maps start with a central idea or topic, 
represented at the center of the map. Branches extend outward from the central topic, representing subtopics 
or related concepts. These branches can further split into sub-branches, forming a hierarchical and 
interconnected structure. provide summary for each branches and sub-branches Generate mind map for this text : {0} 
""".format(documents)
        agent=''
        response_iter = llm.stream_complete(template)
        for response in response_iter:
           print(response.delta, end="", flush=True)
           agent+=response.delta
        print(agent)
        lines = agent.split('\n')
        topic=''
        text1=''
        count=1
        for line in lines[1:]:
            try:
                if 'central idea' in line.lower() and 'note:' not in line.lower():
                  if len(topic)<=0:
                    topic=line.split(':')[1]
                    topic=topic.lstrip()
            except:
                pass
            if 'mind map' in line.lower() and 'note:' not in line.lower():
              if len(topic)<=0:
                topic=line
            elif count==8:
                pass
            elif 'central idea' not in line.lower() and len(line)>0 and 'subtopics' not in line.lower() and  'related concepts' not in line.lower() and 'branches' not in line.lower() and 'Sub-branches' not in line.lower() and 'note:' not in line.lower(): 
                line=line.replace('*','').replace('+','').replace('-','').replace('(','').replace(')','').replace('{','').replace('}','')
                line1=re.sub(r'^\s*\d+(\.\d+)*\.\s*', ' ', line)
                indentation = len(line) - len(line.lstrip())
                line1 = ' ' * indentation + line1.strip()
                line=line1
                if line not in text1:
                    text1+='\t'+line+'\n'
            count+=1

        print(text1,topic)
        if ':' in lines[8] and "+" not in lines[8] and "-" not in lines[8]  and '*' not in lines[8]:
            if len(topic)<=0:
                topic=lines[8].replace(':','')
        if len(topic)<=0:
            topic=file.filename.split('.')[0]
        graph="""
mindmap
root(({0}))
{1}
        """.format(topic,text1)
        print(graph)
        graphbytes = graph.encode("utf-8")
        base64_bytes = base64.b64encode(graphbytes)
        base64_string = base64_bytes.decode("ascii")
        img_url = "https://mermaid.ink/img/" + base64_string
        #display(Image(url=img_url))
        print(img_url)
        save_path = "mermaid_image.png"
        response = requests.get(img_url)
        print(response)
        if response.status_code == 200:
            with open(save_path, 'wb') as file:
                print('abc')
                file.write(response.content)
        else:
            print('def')
        image_path = BASE_DIR+'/'+save_path 
        return send_file(image_path, mimetype='image/jpeg')

if __name__ == "__main__":
    app.run()

